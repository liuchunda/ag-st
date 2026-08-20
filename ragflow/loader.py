from pathlib import Path
from rich import print
import logging

logger = logging.getLogger(__name__)

# 定义当前加载器支持的文件扩展名集合
SUPPORTED_EXTENSIONS = {
    # PDF 文档扩展名
    ".pdf",
    # Word 新版文档扩展名
    ".docx",
    # Word 旧版文档扩展名
    ".doc",
    # Excel 新版表格扩展名
    ".xlsx",
    # Excel 旧版表格扩展名
    ".xls",
    # PowerPoint 新版演示文稿扩展名
    ".pptx",
    # PowerPoint 旧版演示文稿扩展名
    ".ppt",
    # HTML 网页扩展名
    ".html",
    # HTML 网页简写扩展名
    ".htm",
    # XML 文档扩展名
    ".xml",
    # CSV 表格扩展名
    ".csv",
    # JSON 数据扩展名
    ".json",
    # Markdown 文档扩展名
    ".md",
    # 纯文本扩展名
    ".txt",
    # JSON Lines 文本扩展名
    ".jsonl",
}


class DocumentLoader:
    """多格式非结构化文档的加载器"""

    def load(self, file_path):
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在:{path}")
        if not path.is_file():
            raise ValueError(f"路径不是文件:{path}")
        # 获取文件的扩展名，并转为小写
        ext = path.suffix.lower()
        logger.info("解析文档:%s (type=%s)", path.name, ext or "unknown")
        # 构建扩展名到对应的解析方法的映射表
        extractors = {
            # PDF 文件使用 PDF 解析方法
            ".pdf": self._load_pdf,
            # DOCX 文件使用 Word 解析方法
            ".docx": self._load_docx,
            # DOC 文件同样使用 Word 解析方法
            ".doc": self._load_docx,
            # XLSX 文件使用 Excel 解析方法
            ".xlsx": self._load_excel,
            # XLS 文件同样使用 Excel 解析方法
            ".xls": self._load_excel,
            # PPTX 文件使用 PowerPoint 解析方法
            ".pptx": self._load_pptx,
            # PPT 文件同样使用 PowerPoint 解析方法
            ".ppt": self._load_pptx,
            # HTML 文件使用 HTML 解析方法
            ".html": self._load_html,
            # HTM 文件同样使用 HTML 解析方法
            ".htm": self._load_html,
            # XML 文件使用 XML 解析方法
            ".xml": self._load_xml,
            # CSV 文件使用 CSV 解析方法
            ".csv": self._load_csv,
            # JSON 文件使用 JSON 解析方法
            ".json": self._load_json,
            # Markdown 文件使用纯文本解析方法
            ".md": self._load_plain,
            # TXT 文件使用纯文本解析方法
            ".txt": self._load_plain,
            # JSONL 文件使用纯文本解析方法
            ".jsonl": self._load_plain,
        }
        # 根据扩展名找到对应解析函数
        extractor = extractors.get(ext)
        if extractor is None:
            raise ValueError(
                f"不支持的文件类型:{ext}:{','.join(sorted(SUPPORTED_EXTENSIONS))}"
            )
        # 调用对应解析函数提取文档的原始文本
        text = extractor(path)
        # 对提取的原始文本进行清洗和规范化处理
        text = self._normalize(text)
        if not text:
            logger.warning("文档内容为空:%s", path)
        else:
            logger.info("解析完成：%s,字符数=%d", path.name, len(text))
        return text

    @staticmethod
    def _load_pdf(path):
        import fitz  # PyMuPDF

        with fitz.open(path) as pdf:
            return "\n".join(page.get_text("text") for page in pdf)  # type: ignore

    @staticmethod
    def _load_docx(path):
        from docx import Document  # python-docx

        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)

    @staticmethod
    def _load_excel(path):
        import openpyxl  # openpyxl

        # 工作簿 指定整个excel文件
        wb = openpyxl.load_workbook(str(path), data_only=True)
        try:
            rows = []
            # 工作表
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append("\t".join(cells))
            return "\n".join(rows)
        finally:
            wb.close()

    @staticmethod
    def _load_pptx(path):
        from pptx import Presentation  # python-pptx

        ppt = Presentation(str(path))
        # 用于收集全部幻灯片的文本的列表
        texts = []
        for i, slide in enumerate(ppt.slides, start=1):
            slide_texts = [
                shape.text.strip()  # type: ignore
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape.text.strip()  # type: ignore
            ]
            if slide_texts:
                texts.append(f"[Slide {i}]")
                texts.extends(slide_texts)  # type: ignore
        return "\n".join(texts)

    @staticmethod
    def _load_html(path):
        from bs4 import BeautifulSoup  # BeautifulSoup

        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        # 遍历并移除脚本 样式  以及noscript等噪声标签
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)

    @staticmethod
    def _normalize(text):
        """清先多余空白，降低噪声对分块与检索的响应"""
        # 按行拆分，去掉每行的收尾空白
        lines = [line.strip() for line in text.splitlines()]
        # 过滤掉清洗后变成空字符串的行
        lines = [line for line in lines if line]
        return "\n".join(lines)

    @staticmethod
    def _load_plain(path):
        return path.read_text(encoding="utf-8", errors="ignore")

    @staticmethod
    def _load_csv(path):
        import csv

        with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f)
            return "\n".join(",".join(row) for row in reader)

    @staticmethod
    def _load_json(path):
        import json

        data = json.loads(path.read_text(encoding="utf-8"))
        return json.dumps(data, ensure_ascii=False, indent=2)

    @staticmethod
    def _load_xml(path):
        from lxml import etree  # type: ignore

        # 解析XML文件并获取根节点
        root = etree.parse(str(path)).getroot()
        return " ".join(t.strip() for t in root.itertext() if t and t.strip())
