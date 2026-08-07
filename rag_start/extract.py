"""各类文件的文本提取工具，供 6save.py 调用。"""

from __future__ import annotations

import csv
import json
import xml.etree.ElementTree as ET
from pathlib import Path

import pymupdf
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_pdf_text(file_path: str) -> str:
    """从 PDF 提取文本。"""
    doc = pymupdf.open(file_path)
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


def extract_text_from_word(file_path: str) -> str:
    """从 Word 文档提取段落文本。"""
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def extract_text_from_excel(file_path: str) -> str:
    """从 Excel 提取单元格文本。"""
    wb = load_workbook(file_path, data_only=True, read_only=True)
    lines: list[str] = []
    try:
        for sheet in wb.worksheets:
            lines.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if cell is None else str(cell) for cell in row]
                if any(cell.strip() for cell in cells):
                    lines.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(lines)


def extract_ppt_text(file_path: str) -> str:
    """从 PPT 提取文本。"""
    prs = Presentation(file_path)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def extract_text_from_html(file_path: str) -> str:
    """从 HTML 提取可见文本。"""
    html = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_xml_text(file_path: str) -> str:
    """从 XML 提取文本节点。"""
    tree = ET.parse(file_path)
    texts = [text.strip() for text in tree.getroot().itertext() if text and text.strip()]
    return "\n".join(texts)


def read_csv_to_text(file_path: str) -> str:
    """将 CSV 读成文本。"""
    with open(file_path, encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        return "\n".join("\t".join(row) for row in reader)


def extract_text_from_json(file_path: str) -> str:
    """将 JSON 格式化为可读文本。"""
    data = json.loads(Path(file_path).read_text(encoding="utf-8"))
    return json.dumps(data, ensure_ascii=False, indent=2)


def read_text_file(file_path: str) -> str:
    """读取纯文本 / Markdown / JSONL。"""
    return Path(file_path).read_text(encoding="utf-8", errors="ignore")
